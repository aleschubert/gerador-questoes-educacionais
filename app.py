import streamlit as st
import random
from docx import Document
from io import BytesIO
from fpdf import FPDF
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import re
import base64
import unicodedata

# --- Utilitárias de limpeza preservando acentos ---
def remover_controle(texto):
    # remove caracteres de controle
    return "".join(ch for ch in texto if unicodedata.category(ch)[0] != "C")

def normalizar_texto(texto):
    if texto is None:
        return ""
    # Normaliza unicode (preserva acentos)
    texto = unicodedata.normalize("NFC", texto)
    texto = remover_controle(texto)
    # Troca múltiplos espaços/quebras por um único espaço
    texto = re.sub(r'\s+', ' ', texto)
    return texto.strip()

# --- Extração de texto ---
def extrair_texto_pdf(file):
    try:
        reader = PyPDF2.PdfReader(file)
        texto = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                texto += page_text + " "
        return normalizar_texto(texto)
    except Exception:
        return ""

def extrair_texto_word(file):
    try:
        doc = DocxDocument(file)
        texto = ""
        for para in doc.paragraphs:
            if para.text and para.text.strip():
                texto += para.text + " "
        return normalizar_texto(texto)
    except Exception:
        return ""

def extrair_texto_pptx(file):
    try:
        prs = Presentation(file)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texto += shape.text + " "
        return normalizar_texto(texto)
    except Exception:
        return ""

# --- Separador de sentenças robusto (preserva sentenças inteiras) ---
def separar_sentencas(texto):
    texto = texto.strip()
    if not texto:
        return []
    # split por fim de frase (.!?), mantendo o delimitador na sentença
    partes = re.split(r'(?<=[\.\?\!])\s+', texto)
    # filtrar sentenças muito curtas
    sentencas = [p.strip() for p in partes if len(p.strip()) >= 25]
    return sentencas

# --- Geração de questão (sem GPT) ---
def gerar_questao_inventada(texto_base):
    sentencas = separar_sentencas(texto_base)
    if not sentencas:
        contexto = "Texto insuficiente para gerar questão."
    else:
        contexto = random.choice(sentencas)

    # Enunciado claro usando a própria frase como referência
    enunciado = f"Com base no trecho acima, assinale a alternativa que corresponde corretamente à informação apresentada."

    # Alternativa correta: a frase completa (contexto)
    correta = contexto

    # Alternativas incorretas inventadas (sem relação direta)
    incorretas = [
        "Afirmação falsa para confundir o leitor.",
        "Fato não relacionado ao conteúdo apresentado.",
        "Esta alternativa não corresponde ao texto.",
        "Informação incorreta sobre o assunto."
    ]
    alternativas = [correta] + random.sample(incorretas, 4)
    random.shuffle(alternativas)

    return {
        "contexto": contexto,
        "enunciado": enunciado,
        "alternativas": alternativas,
        "correta": correta
    }

# --- Exportadores ---
def gerar_word(questoes):
    doc = Document()
    doc.add_heading("Prova Gerada", level=1)
    for i, q in enumerate(questoes, 1):
        doc.add_paragraph(f"Questão {i}")
        doc.add_paragraph(q.get("contexto", ""))
        doc.add_paragraph(q.get("enunciado", ""))
        for alt in q.get("alternativas", []):
            doc.add_paragraph(f"- {alt}")
        doc.add_paragraph(f"Gabarito: {q.get('correta','')}")
        doc.add_paragraph("")
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def gerar_pdf(questoes):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for i, q in enumerate(questoes, 1):
        # garantir que texto passado para PDF seja string limpa
        contexto = str(q.get("contexto", ""))
        enunciado = str(q.get("enunciado", ""))
        pdf.multi_cell(0, 8, f"Questão {i}")
        pdf.multi_cell(0, 8, contexto)
        pdf.multi_cell(0, 8, enunciado)
        for alt in q.get("alternativas", []):
            pdf.multi_cell(0, 8, f"- {str(alt)}")
        pdf.multi_cell(0, 8, f"Gabarito: {str(q.get('correta',''))}")
        pdf.ln(5)
    buffer = BytesIO()
    pdf.output(buffer)
    buffer.seek(0)
    return buffer

# --- Streamlit UI ---
st.set_page_config(page_title="Gerador de Questões (versão limpa)", layout="centered")
st.title("Gerador de Questões – Versão Gratuita (enunciados claros)")
st.write("Cole o texto ou envie um arquivo (PDF, Word, PowerPoint). A alternativa correta vem do texto; as incorretas são inventadas.")

texto = st.text_area("Cole o texto aqui (ou envie um arquivo):", height=200)

uploaded_file = st.file_uploader("Ou envie um arquivo:", type=['pdf','docx','pptx'])
if uploaded_file is not None:
    if uploaded_file.name.lower().endswith('.pdf'):
        texto = extrair_texto_pdf(uploaded_file)
    elif uploaded_file.name.lower().endswith('.docx'):
        texto = extrair_texto_word(uploaded_file)
    elif uploaded_file.name.lower().endswith('.pptx'):
        texto = extrair_texto_pptx(uploaded_file)

quantidade = st.number_input("Quantas questões deseja gerar?", min_value=1, max_value=50, value=5)

if st.button("🧠 Gerar questões"):
    if not texto or not texto.strip():
        st.warning("Insira um texto válido ou envie um arquivo com conteúdo.")
    else:
        questoes = []
        sentencas = separar_sentencas(texto)
        if len(sentencas) == 0:
            st.warning("O texto fornecido não possui sentenças completas suficientes. Tente outro documento.")
        else:
            for i in range(quantidade):
                q = gerar_questao_inventada(texto)
                questoes.append(q)
                st.subheader(f"Questão {i+1}")
                st.write(q["contexto"])
                st.markdown(f"**{q['enunciado']}**")
                for alt in q["alternativas"]:
                    st.write(f"- {alt}")
                st.success(f"✔️ Alternativa correta (texto integral):")
                st.write(q["correta"])
                st.divider()

            # oferecer downloads
            buf_word = gerar_word(questoes)
            b64w = base64.b64encode(buf_word.read()).decode()
            hrefw = f'<a href="data:application/octet-stream;base64,{b64w}" download="Prova.docx">💾 Baixar em Word</a>'
            st.markdown(hrefw, unsafe_allow_html=True)

            buf_pdf = gerar_pdf(questoes)
            b64p = base64.b64encode(buf_pdf.read()).decode()
            hrefp = f'<a href="data:application/octet-stream;base64,{b64p}" download="Prova.pdf">💾 Baixar em PDF</a>'
            st.markdown(hrefp, unsafe_allow_html=True)
            
